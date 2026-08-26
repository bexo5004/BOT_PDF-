"""
file_engine.py - محرك معالجة الملفات المتكامل
يدعم: PDF, Word, Excel, PowerPoint, PPSX, صور, نصوص, أرشيف
"""

import os
import io
import time
import shutil
import zipfile
from pathlib import Path
from typing import List, Tuple, Optional

from pypdf import PdfReader, PdfWriter
import fitz
from PIL import Image
from config import Config
from utils import logger, safe_remove, format_size

# ============================================================
# المكتبات الاختيارية
# ============================================================
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False


class FileEngine:
    """محرك معالجة الملفات المتكامل"""

    # ============================================================
    # 1. تحديد نوع الملف
    # ============================================================
    
    @staticmethod
    def get_file_type(file_path: str) -> str:
        """تحديد نوع الملف"""
        ext = Path(file_path).suffix.lower()
        
        # مستندات
        if ext == '.pdf':
            return 'pdf'
        if ext in ['.doc', '.docx', '.docm']:
            return 'word'
        if ext in ['.xls', '.xlsx', '.xlsm', '.xlsb', '.csv']:
            return 'excel'
        if ext in ['.ppt', '.pptx', '.pptm', '.ppsx', '.pps']:
            return 'powerpoint'
        if ext in ['.txt', '.rtf', '.md', '.markdown']:
            return 'text'
        if ext in ['.odt', '.ods', '.odp']:
            return 'openoffice'
        if ext in ['.html', '.htm', '.xml', '.json']:
            return 'web'
        if ext in ['.epub', '.mobi', '.fb2']:
            return 'ebook'
        if ext == '.tex':
            return 'latex'
        
        # صور
        if ext in ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff', '.tif', '.gif', '.svg', '.heic', '.ico']:
            return 'image'
        
        # أرشيف
        if ext in ['.zip', '.rar', '.7z', '.gz', '.tar', '.tgz']:
            return 'archive'
        
        return 'unknown'

    # ============================================================
    # 2. تحويل الملفات إلى PDF
    # ============================================================
    
    @staticmethod
    def convert_to_pdf(file_path: str) -> str:
        """تحويل أي ملف إلى PDF"""
        file_type = FileEngine.get_file_type(file_path)
        
        if file_type == 'pdf':
            return file_path
        if file_type == 'image':
            return FileEngine._image_to_pdf(file_path)
        if file_type == 'powerpoint':
            return FileEngine._powerpoint_to_pdf(file_path)
        if file_type == 'word':
            return FileEngine._word_to_pdf(file_path)
        if file_type == 'excel':
            return FileEngine._excel_to_pdf(file_path)
        if file_type in ['text', 'openoffice', 'web', 'latex']:
            return FileEngine._text_to_pdf(file_path)
        if file_type == 'ebook':
            return FileEngine._ebook_to_pdf(file_path)
        
        # محاولة معالجة الملفات غير المعروفة كنصوص
        try:
            return FileEngine._text_to_pdf(file_path)
        except:
            raise ValueError(f"نوع الملف {file_type} غير مدعوم")

    # ============================================================
    # 3. تحويل النصوص إلى PDF
    # ============================================================
    
    @staticmethod
    def text_to_pdf(text: str, title: str = "نص") -> str:
        """تحويل نص إلى PDF مع دعم العربية"""
        try:
            out_path = Path(Config.TEMP_DIR) / f"text_{os.urandom(4).hex()}.pdf"
            
            if not REPORTLAB_AVAILABLE:
                raise RuntimeError("reportlab غير مثبت")
            
            # البحث عن خط يدعم العربية
            font_name = "Helvetica"
            font_paths = [
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/noto/NotoSansArabic-Regular.ttf",
            ]
            
            for path in font_paths:
                if os.path.exists(path):
                    try:
                        pdfmetrics.registerFont(TTFont("ArabicFont", path))
                        font_name = "ArabicFont"
                        break
                    except:
                        continue
            
            c = canvas.Canvas(str(out_path), pagesize=A4)
            width, height = A4
            
            margin = 72
            x = margin
            y = height - margin
            
            # العنوان
            c.setFont(font_name, 20)
            c.setFillColor(colors.darkblue)
            c.drawString(x, y, f"📄 {title}")
            y -= 40
            
            # خط فاصل
            c.setStrokeColor(colors.black)
            c.line(x, y, width - margin, y)
            y -= 30
            
            # النص
            c.setFont(font_name, 12)
            c.setFillColor(colors.black)
            
            lines = text.split('\n')
            for line in lines:
                if y < 50:
                    c.showPage()
                    y = height - margin
                    c.setFont(font_name, 12)
                    c.setFillColor(colors.black)
                
                if line.strip():
                    clean_line = line.strip()
                    if len(clean_line) > 80:
                        words = clean_line.split(' ')
                        current_line = ""
                        for word in words:
                            if len(current_line) + len(word) + 1 < 80:
                                current_line += word + " "
                            else:
                                c.drawString(x, y, current_line.strip())
                                y -= 18
                                current_line = word + " "
                        if current_line:
                            c.drawString(x, y, current_line.strip())
                            y -= 18
                    else:
                        c.drawString(x, y, clean_line)
                        y -= 18
            
            # التذييل مع حقوق البوت
            y -= 20
            c.setFont(font_name, 10)
            c.setFillColor(colors.grey)
            c.drawString(x, y, f"PDF Bot • @BEXO50 • {time.strftime('%Y-%m-%d %H:%M')}")
            
            c.save()
            
            logger.info(f"✅ تم تحويل النص إلى PDF: {Path(out_path).name}")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل تحويل النص إلى PDF: {str(e)}")
    
    @staticmethod
    def _text_to_pdf(file_path: str) -> str:
        """تحويل ملف نصي إلى PDF"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return FileEngine.text_to_pdf(content, Path(file_path).stem)
        except Exception as e:
            raise RuntimeError(f"فشل تحويل الملف النصي: {str(e)}")

    # ============================================================
    # 4. تحويل الصور إلى PDF
    # ============================================================
    
    @staticmethod
    def _image_to_pdf(file_path: str) -> str:
        """تحويل صورة إلى PDF"""
        try:
            with Image.open(file_path) as img:
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                out_path = Path(Config.TEMP_DIR) / f"img_{os.urandom(4).hex()}.pdf"
                img.save(str(out_path), "PDF", resolution=100.0)
                logger.info(f"✅ تحويل صورة: {Path(file_path).name}")
                return str(out_path)
        except Exception as e:
            raise RuntimeError(f"فشل تحويل الصورة: {e}")
    
    @staticmethod
    def images_to_pdf(image_paths: List[str]) -> str:
        """تحويل مجموعة صور إلى PDF"""
        if not image_paths:
            raise ValueError("لا توجد صور")
        
        images = []
        for path in image_paths:
            try:
                img = Image.open(path)
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                images.append(img)
            except Exception as e:
                logger.error(f"❌ فشل تحميل الصورة: {e}")
        
        if not images:
            raise ValueError("لا توجد صور صالحة")
        
        out_path = Path(Config.TEMP_DIR) / f"images_{os.urandom(4).hex()}.pdf"
        if len(images) == 1:
            images[0].save(str(out_path), "PDF")
        else:
            images[0].save(str(out_path), "PDF", save_all=True, append_images=images[1:])
        
        for img in images:
            img.close()
        
        logger.info(f"✅ تم تحويل {len(images)} صورة إلى PDF")
        return str(out_path)

    # ============================================================
    # 5. تحويل PowerPoint إلى PDF
    # ============================================================
    
    @staticmethod
    def _powerpoint_to_pdf(file_path: str) -> str:
        """تحويل PowerPoint إلى PDF"""
        out_path = Path(Config.TEMP_DIR) / f"ppt_{os.urandom(4).hex()}.pdf"
        
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx غير مثبت")
        
        try:
            prs = Presentation(file_path)
            slides_text = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slides_text.append(text)
            
            if not slides_text:
                raise ValueError("لا يوجد نص في العرض")
            
            if not REPORTLAB_AVAILABLE:
                raise RuntimeError("reportlab غير مثبت")
            
            c = canvas.Canvas(str(out_path), pagesize=A4)
            y = 750
            for text in slides_text[:100]:
                if y < 50:
                    c.showPage()
                    y = 750
                c.drawString(50, y, text[:200])
                y -= 20
            c.save()
            
            logger.info(f"✅ تحويل PowerPoint: {Path(file_path).name}")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل تحويل PowerPoint: {str(e)}")

    # ============================================================
    # 6. تحويل Word إلى PDF
    # ============================================================
    
    @staticmethod
    def _word_to_pdf(file_path: str) -> str:
        """تحويل Word إلى PDF"""
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx غير مثبت")
        
        try:
            doc = Document(file_path)
            out_path = Path(Config.TEMP_DIR) / f"word_{os.urandom(4).hex()}.pdf"
            
            if not REPORTLAB_AVAILABLE:
                raise RuntimeError("reportlab غير مثبت")
            
            c = canvas.Canvas(str(out_path), pagesize=A4)
            y = 750
            
            for para in doc.paragraphs[:100]:
                if para.text.strip():
                    if y < 50:
                        c.showPage()
                        y = 750
                    c.drawString(50, y, para.text[:200])
                    y -= 20
            
            c.save()
            logger.info(f"✅ تحويل Word: {Path(file_path).name}")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل تحويل Word: {str(e)}")

    # ============================================================
    # 7. تحويل Excel إلى PDF
    # ============================================================
    
    @staticmethod
    def _excel_to_pdf(file_path: str) -> str:
        """تحويل Excel إلى PDF"""
        if not OPENPYXL_AVAILABLE:
            raise RuntimeError("openpyxl غير مثبت")
        
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            out_path = Path(Config.TEMP_DIR) / f"excel_{os.urandom(4).hex()}.pdf"
            
            if not REPORTLAB_AVAILABLE:
                raise RuntimeError("reportlab غير مثبت")
            
            c = canvas.Canvas(str(out_path), pagesize=A4)
            y = 750
            
            for sheet_name in wb.sheetnames[:3]:
                sheet = wb[sheet_name]
                c.drawString(50, y, f"📊 {sheet_name}")
                y -= 30
                
                for row in sheet.iter_rows(values=True):
                    if y < 50:
                        c.showPage()
                        y = 750
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        c.drawString(50, y, row_text[:200])
                        y -= 15
                y -= 20
            
            c.save()
            logger.info(f"✅ تحويل Excel: {Path(file_path).name}")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل تحويل Excel: {str(e)}")

    # ============================================================
    # 8. تحويل الكتب الإلكترونية إلى PDF
    # ============================================================
    
    @staticmethod
    def _ebook_to_pdf(file_path: str) -> str:
        """تحويل كتاب إلكتروني إلى PDF"""
        out_path = Path(Config.TEMP_DIR) / f"ebook_{os.urandom(4).hex()}.pdf"
        
        if not REPORTLAB_AVAILABLE:
            raise RuntimeError("reportlab غير مثبت")
        
        c = canvas.Canvas(str(out_path), pagesize=A4)
        c.drawString(50, 750, f"📖 {Path(file_path).stem}")
        c.drawString(50, 730, f"الحجم: {format_size(os.path.getsize(file_path))}")
        c.drawString(50, 710, f"النوع: {Path(file_path).suffix}")
        c.save()
        
        return str(out_path)

    # ============================================================
    # 9. دمج المستندات
    # ============================================================
    
    @staticmethod
    def merge_documents(file_paths: List[str]) -> str:
        """دمج مستندات متعددة في PDF واحد"""
        if not file_paths:
            raise ValueError("لا توجد ملفات")
        
        logger.info(f"🔀 دمج {len(file_paths)} ملف")
        
        pdf_files = []
        errors = []
        
        for path in file_paths:
            try:
                pdf_path = FileEngine.convert_to_pdf(path)
                pdf_files.append(pdf_path)
                logger.info(f"✅ تم تحويل: {Path(path).name}")
            except Exception as e:
                errors.append(f"{Path(path).name}: {str(e)}")
                logger.error(f"❌ فشل {Path(path).name}: {e}")
        
        if not pdf_files:
            error_details = "\n".join(errors[:5])
            raise RuntimeError(f"لا توجد ملفات صالحة للدمج\n\nالأخطاء:\n{error_details}")
        
        try:
            writer = PdfWriter()
            for path in pdf_files:
                reader = PdfReader(path)
                for page in reader.pages:
                    writer.add_page(page)
            
            out_path = Path(Config.TEMP_DIR) / f"merged_{os.urandom(4).hex()}.pdf"
            with open(out_path, "wb") as f:
                writer.write(f)
            
            for path in pdf_files:
                if path not in file_paths:
                    safe_remove(path)
            
            logger.info(f"✅ تم دمج {len(pdf_files)} ملف")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل الدمج: {str(e)}")

    # ============================================================
    # 10. استخراج الصور من PDF
    # ============================================================
    
    @staticmethod
    def extract_images_from_pdf(pdf_path: str, dpi: int = 150) -> Tuple[bytes, str]:
        """استخراج الصور من PDF"""
        doc = None
        try:
            doc = fitz.open(pdf_path)
            
            if len(doc) == 0:
                raise ValueError("الملف فارغ")
            
            if len(doc) == 1:
                page = doc[0]
                pix = page.get_pixmap(dpi=dpi)
                img_data = pix.tobytes("jpeg")
                doc.close()
                return img_data, "صفحة_1.jpg"
            
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                for i, page in enumerate(doc, 1):
                    pix = page.get_pixmap(dpi=dpi)
                    img_data = pix.tobytes("jpeg")
                    zf.writestr(f"صفحة_{i}.jpg", img_data)
            
            doc.close()
            return buf.getvalue(), "صور_مستخرجة.zip"
            
        except Exception as e:
            if doc:
                try:
                    doc.close()
                except:
                    pass
            raise RuntimeError(f"فشل استخراج الصور: {str(e)}")

    # ============================================================
    # 11. تقسيم PDF
    # ============================================================
    
    @staticmethod
    def split_pdf(pdf_path: str, page_ranges: str) -> str:
        """تقسيم PDF واستخراج صفحات محددة"""
        try:
            pages = set()
            for part in page_ranges.replace(" ", "").split(","):
                if not part:
                    continue
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    pages.update(range(start, end + 1))
                else:
                    pages.add(int(part))
            
            if not pages:
                raise ValueError("لم يتم تحديد أي صفحات")
            
            reader = PdfReader(pdf_path)
            total_pages = len(reader.pages)
            valid_pages = sorted([p for p in pages if 1 <= p <= total_pages])
            
            if not valid_pages:
                raise ValueError(f"لا توجد صفحات صالحة (1-{total_pages})")
            
            writer = PdfWriter()
            for page_num in valid_pages:
                writer.add_page(reader.pages[page_num - 1])
            
            out_path = Path(Config.TEMP_DIR) / f"split_{os.urandom(4).hex()}.pdf"
            with open(out_path, "wb") as f:
                writer.write(f)
            
            logger.info(f"✅ تم تقسيم PDF: استخراج {len(valid_pages)} صفحة")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل تقسيم PDF: {str(e)}")

    # ============================================================
    # 12. حذف صفحات
    # ============================================================
    
    @staticmethod
    def delete_pages(pdf_path: str, pages_to_delete: List[int]) -> str:
        """حذف صفحات محددة"""
        try:
            reader = PdfReader(pdf_path)
            total_pages = len(reader.pages)
            
            if total_pages == 0:
                raise ValueError("الملف فارغ")
            
            delete_set = set(pages_to_delete)
            remaining = [i for i in range(1, total_pages + 1) if i not in delete_set]
            
            if not remaining:
                raise ValueError("لا يمكن حذف جميع الصفحات")
            
            writer = PdfWriter()
            for i in range(total_pages):
                if (i + 1) not in delete_set:
                    writer.add_page(reader.pages[i])
            
            out_path = Path(Config.TEMP_DIR) / f"deleted_{os.urandom(4).hex()}.pdf"
            with open(out_path, "wb") as f:
                writer.write(f)
            
            logger.info(f"✅ تم حذف {len(pages_to_delete)} صفحة")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل حذف الصفحات: {str(e)}")

    # ============================================================
    # 13. ترقيم الصفحات
    # ============================================================
    
    @staticmethod
    def add_page_numbers(pdf_path: str) -> str:
        """إضافة أرقام الصفحات"""
        try:
            doc = fitz.open(pdf_path)
            
            for i, page in enumerate(doc, 1):
                rect = fitz.Rect(
                    page.rect.width * 0.45,
                    page.rect.height - 40,
                    page.rect.width * 0.55,
                    page.rect.height - 10
                )
                page.insert_textbox(rect, str(i), fontsize=14, color=(0, 0, 0))
            
            out_path = Path(Config.TEMP_DIR) / f"numbered_{os.urandom(4).hex()}.pdf"
            doc.save(str(out_path))
            doc.close()
            
            logger.info(f"✅ تمت إضافة أرقام الصفحات")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل إضافة الأرقام: {str(e)}")

    # ============================================================
    # 14. ضغط PDF
    # ============================================================
    
    @staticmethod
    def compress_pdf(pdf_path: str) -> Tuple[str, int, int]:
        """ضغط PDF"""
        before = os.path.getsize(pdf_path)
        
        try:
            doc = fitz.open(pdf_path)
            out_path = Path(Config.TEMP_DIR) / f"compressed_{os.urandom(4).hex()}.pdf"
            doc.save(str(out_path), garbage=4, deflate=True, clean=True)
            doc.close()
            
            after = os.path.getsize(out_path)
            
            if after >= before:
                safe_remove(out_path)
                return pdf_path, before, before
            
            logger.info(f"✅ تم ضغط PDF: {format_size(before)} → {format_size(after)}")
            return str(out_path), before, after
            
        except Exception as e:
            raise RuntimeError(f"فشل ضغط PDF: {str(e)}")

    # ============================================================
    # 15. تشفير PDF
    # ============================================================
    
    @staticmethod
    def encrypt_pdf(pdf_path: str, password: str) -> str:
        """تشفير PDF بكلمة مرور"""
        if not password or len(password) < 4:
            raise ValueError("كلمة المرور 4 أحرف على الأقل")
        
        try:
            reader = PdfReader(pdf_path)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(password)
            
            out_path = Path(Config.TEMP_DIR) / f"encrypted_{os.urandom(4).hex()}.pdf"
            with open(out_path, "wb") as f:
                writer.write(f)
            
            logger.info(f"✅ تم تشفير PDF")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل التشفير: {str(e)}")

    # ============================================================
    # 16. إزالة الحماية
    # ============================================================
    
    @staticmethod
    def remove_password(pdf_path: str) -> str:
        """إزالة الحماية من PDF"""
        try:
            reader = PdfReader(pdf_path)
            if not reader.is_encrypted:
                raise ValueError("الملف غير مشفر")
            
            try:
                reader.decrypt('')
            except:
                raise ValueError("لا يمكن إزالة الحماية بدون كلمة المرور")
            
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            
            out_path = Path(Config.TEMP_DIR) / f"unlocked_{os.urandom(4).hex()}.pdf"
            with open(out_path, "wb") as f:
                writer.write(f)
            
            logger.info(f"✅ تم إزالة الحماية")
            return str(out_path)
            
        except Exception as e:
            raise RuntimeError(f"فشل إزالة الحماية: {str(e)}")

    # ============================================================
    # 17. فك ضغط الأرشيف
    # ============================================================
    
    @staticmethod
    def extract_archive(file_path: str) -> List[str]:
        """فك ضغط الملفات المضغوطة"""
        ext = Path(file_path).suffix.lower()
        output_dir = Path(Config.TEMP_DIR) / f"extract_{os.urandom(4).hex()}"
        output_dir.mkdir(exist_ok=True)
        
        try:
            if ext in ['.zip', '.zipx']:
                shutil.unpack_archive(file_path, output_dir, 'zip')
            elif ext == '.rar':
                try:
                    import rarfile
                    with rarfile.RarFile(file_path) as rf:
                        rf.extractall(output_dir)
                except ImportError:
                    raise ValueError("rarfile غير مثبت")
            elif ext == '.7z':
                try:
                    import py7zr
                    with py7zr.SevenZipFile(file_path, mode='r') as z:
                        z.extractall(output_dir)
                except ImportError:
                    raise ValueError("py7zr غير مثبت")
            elif ext in ['.gz', '.tar', '.tgz', '.tar.gz', '.bz2', '.xz']:
                shutil.unpack_archive(file_path, output_dir)
            else:
                raise ValueError(f"نوع الأرشيف {ext} غير مدعوم")
            
            extracted = []
            for root, dirs, files in os.walk(output_dir):
                for file in files:
                    extracted.append(os.path.join(root, file))
            
            return extracted
            
        except Exception as e:
            shutil.rmtree(output_dir, ignore_errors=True)
            raise RuntimeError(f"فشل فك الضغط: {str(e)}")
